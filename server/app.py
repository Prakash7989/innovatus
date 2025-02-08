from serverless_wsgi import handle_request
import io
import os
import re
import warnings
from threading import Thread
from flask import Flask, jsonify, request, current_app
from flask_cors import CORS
from werkzeug.utils import secure_filename
import pdfplumber
from pymongo import MongoClient
from bson import ObjectId, Binary
from docx import Document
from pptx import Presentation
from transformers import BartForConditionalGeneration, BartTokenizer, pipeline
import torch
from datetime import datetime
from currentsapi import CurrentsAPI
from bson.errors import InvalidId


# Initialize Flask app
app = Flask(__name__)
CORS(app)

# Suppress warnings
warnings.filterwarnings("ignore")

# Configure device
device = "cuda" if torch.cuda.is_available() else "cpu"

# Load ML models
try:
    summary_model_name = "facebook/bart-large-cnn"
    tokenizer = BartTokenizer.from_pretrained(summary_model_name)
    summary_model = BartForConditionalGeneration.from_pretrained(summary_model_name).to(device)
    if device == "cuda":
        summary_model = summary_model.half()
    classifier = pipeline(
        "zero-shot-classification",
        model="facebook/bart-large-mnli",
        device=0 if device == "cuda" else -1,
        batch_size=8
    )
    print("✅ Models loaded successfully!")
except Exception as e:
    print(f"❌ Error loading models: {str(e)}")
    exit(1)

# MongoDB configuration
client = MongoClient(os.getenv("mongodb+srv://jeevaniitkgp22:Jeevan2025@cluster0.dzvo6.mongodb.net"))
file_db = client["file_storage_db"]
file_collection = file_db["files"]

# Separate databases and collections
news_db = client["news_db"]
news_collection = news_db["news"]

# File upload configuration
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pdf', 'ppt', 'pptx', 'doc', 'docx'}
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024

# Classification categories
CATEGORIES = [
    "Software Development", "Cybersecurity", "Cloud Computing",
    "Artificial Intelligence", "Machine Learning", "Deep Learning",
    "Business Strategy", "Stock Market", "Cryptocurrency",
    "Physics", "Chemistry", "Biology", "Medicine",
    "Education", "Online Learning", "Mental Health",
    "Journalism", "Social Media", "Marketing",
    "Sports", "E-Sports", "History", "Politics",
    "Renewable Energy", "Climate Change", "Environment",
    "Food & Cuisine", "Travel", "Personal Development"
]

# Helper functions
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def clean_text(text):
    if not text:
        return ""
    return re.sub(r'\s+', ' ', text).strip()

def extract_text_from_docx(file_bytes):
    try:
        doc = Document(io.BytesIO(file_bytes))
        return clean_text('\n'.join([para.text for para in doc.paragraphs]))
    except Exception as e:
        return f"Error extracting DOCX: {str(e)}"

def extract_text_from_pptx(file_bytes):
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        return clean_text('\n'.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]))
    except Exception as e:
        return f"Error extracting PPT/PPTX: {str(e)}"

def generate_summary(text):
    """Generate a summary using BART."""
    inputs = tokenizer(text, return_tensors="pt", truncation=True, max_length=1024)
    inputs = {key: value.to(device) for key, value in inputs.items()}
    
    summary_ids = summary_model.generate(
        **inputs, max_length=200, num_beams=4, early_stopping=True
    )
    
    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    return summary

def classify_text(text):
    """Classify text into top 3 categories using zero-shot classification."""
    result = classifier(text, CATEGORIES, multi_label=True)
    top_categories = sorted(zip(result["labels"], result["scores"]), key=lambda x: x[1], reverse=True)[:3]
    
    return [{"category": cat, "score": round(score, 2)} for cat, score in top_categories]

@app.route('/upload', methods=['POST'])

def upload_file():
    
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files['file']
    if not file or file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    if not allowed_file(file.filename):
        return jsonify({"error": "Invalid file type"}), 400

    filename = secure_filename(file.filename)
    file_bytes = file.read()

    file_data = {
        "filename": filename,
        "file_type": file.mimetype,
        "content": Binary(file_bytes),  
        "status": "pending",
        "upload_date": datetime.utcnow().isoformat()
    }

    inserted_id = file_collection.insert_one(file_data).inserted_id

    # Start background processing
    Thread(target=process_file_background, args=(str(inserted_id),)).start()

    return jsonify({"message": "File uploaded successfully", "file_id": str(inserted_id)}), 200

def process_file_background(file_id):
    with app.app_context():
        try:
            file_data = file_collection.find_one({"_id": ObjectId(file_id)})
            if not file_data:
                return

            # Extract text from binary content
            file_bytes = file_data['content']
            text = ""
            if file_data['filename'].endswith(".pdf"):
                with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                    text = " ".join([page.extract_text() or "" for page in pdf.pages])
            elif file_data['filename'].endswith(".docx"):
                text = extract_text_from_docx(file_bytes)
            elif file_data['filename'].endswith((".ppt", ".pptx")):
                text = extract_text_from_pptx(file_bytes)

            text = clean_text(text)

            # Generate summary and categories
            summary = generate_summary(text)
            categories = classify_text(text)

            # Update database
            file_collection.update_one(
                {"_id": ObjectId(file_id)},
                {"$set": {
                    "text": text,
                    "summary": summary,
                    "categories": categories,
                    "status": "processed"
                }}
            )
        except Exception as e:
            print(f"Error processing {file_id}: {str(e)}")
            file_collection.update_one(
                {"_id": ObjectId(file_id)},
                {"$set": {"status": "error", "error": str(e)}}
            )

@app.route('/get-summary/<file_id>', methods=['GET'])
def get_summary(file_id):
    try:
        object_id = ObjectId(file_id)
    except InvalidId:
        return jsonify({"error": "Invalid file ID"}), 400

    file_data = file_collection.find_one({"_id": object_id})
    if not file_data:
        return jsonify({"error": "File not found"}), 404

    status = file_data.get("status", "pending")
    
    if status == "processed":
        return jsonify({
            "summary": file_data.get("summary", ""),
            "categories": file_data.get("categories", []),
            "originaltext": file_data.get("text", "")
        }), 200
    elif status == "pending":
        return jsonify({"status": "processing", "message": "File is being processed"}), 202
    else:
        return jsonify({
            "error": "Processing error",
            "details": file_data.get("error", "")
        }), 500



@app.route('/files', methods=['GET'])
def get_all_files():
    try:
        files = list(file_collection.find({}, {
            "_id": 1, 
            "filename": 1, 
            "file_type": 1, 
            "status": 1,
            "summary": 1, 
            "categories": 1
        }))
        return jsonify([{
            "id": str(f["_id"]),
            "name": f["filename"],
            "type": f["file_type"],
            "status": f.get("status", "pending"),
            "summary": f.get("summary", ""),
            "categories": f.get("categories", [])
        } for f in files]), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/delete-file/<file_id>', methods=['DELETE'])
def delete_file(file_id):
    try:
        object_id = ObjectId(file_id)
    except InvalidId:
        return jsonify({"error": "Invalid file ID"}), 400

    file_data = file_collection.find_one({"_id": object_id})
    if not file_data:
        return jsonify({"error": "File not found"}), 404

    file_collection.delete_one({"_id": object_id})
    return jsonify({"message": "File deleted successfully"}), 200



currents_api_key = "_NSfqmF4m1zJsae21ns7FØRPOB1aHvhNQVSj5N7c2dnHG7Tx"
if not currents_api_key:
    raise ValueError("CURRENTS_API_KEY environment variable is not set. Please set it before running the application.")
api = CurrentsAPI(api_key=currents_api_key)

@app.route("/fetch-news", methods=["GET"])
def fetch_news():
    try:
        news_data = api.latest_news()
        articles = news_data.get("news", [])

        for article in articles:
            article["stored_at"] = datetime.utcnow()
            news_collection.update_one({"url": article["url"]}, {"$set": article}, upsert=True)

        return jsonify({"message": "News articles fetched and stored successfully!"}), 200
    except Exception as e:
        return jsonify({"error": f"Failed to fetch news. Error: {str(e)}"}), 500

@app.route("/get-news", methods=["GET"])
def get_news():
    news_articles = list(news_collection.find({}, {"_id": 0})) 
    return jsonify(news_articles), 200
# Run Flask app
if __name__ == '__main__':
    app.run(debug=True)

def handler(event, context):
    return handle_request(app, event, context)